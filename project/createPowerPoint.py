from pptx import Presentation
from pptx.util import Inches, Pt
import os

def create_presentation():
    # Create a blank presentation
    prs = Presentation()

    # --- HELPER FUNCTION: Add Title Slide ---
    def add_title_slide(title, subtitle):
        slide_layout = prs.slide_layouts[0] # 0 = Title Slide
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    # --- HELPER FUNCTION: Add Bullet Slide ---
    def add_bullet_slide(title, content_list):
        slide_layout = prs.slide_layouts[1] # 1 = Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        
        tf = slide.shapes.placeholders[1].text_frame
        for i, line in enumerate(content_list):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(20) # Set readable font size

    # --- HELPER FUNCTION: Add Image Slide ---
    def add_image_slide(title, image_path, description_list=[]):
        slide_layout = prs.slide_layouts[5] # 5 = Title Only (we add image manually)
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        
        # Add Image (if it exists)
        if os.path.exists(image_path):
            # Centered image
            left = Inches(1)
            top = Inches(2)
            height = Inches(4.5)
            slide.shapes.add_picture(image_path, left, top, height=height)
        else:
            # Placeholder text if image is missing
            txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(5), Inches(1))
            txBox.text_frame.text = f"[Image '{image_path}' not found. Please run analysis script first.]"

        # Add Description Text below or to the side
        if description_list:
            txBox = slide.shapes.add_textbox(Inches(6), Inches(2), Inches(3.5), Inches(4))
            tf = txBox.text_frame
            tf.word_wrap = True
            for line in description_list:
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(16)
                p.level = 0

    # ==========================================
    # SLIDE 1: Title
    # ==========================================
    add_title_slide("Social Networks Analysis", "EU Email Network Analysis")

    # ==========================================
    # SLIDE 2: Dataset Information
    # ==========================================
    # Summarized points as requested (Source, Time, Nodes, Edges)
    dataset_points = [
        "Source: Large European Research Institution",
        "Timeframe: October 2003 – May 2005 (18 Months)",
        "Structure: Directed Graph (Edge i → j if i sent email to j)",
        " ",
        "--- Key Statistics ---",
        "Nodes (Email Addresses): 265,214",
        "Edges (Sent Emails): 420,045",
        "Avg Clustering Coefficient: 0.067",
        "Reciprocity: Only ~13% of nodes are in the largest Strongly Connected Component (SCC)"
    ]
    add_bullet_slide("Dataset Information", dataset_points)

    # ==========================================
    # SLIDE 3: Senders vs Receivers
    # ==========================================
    add_image_slide(
        "Top Senders vs. Top Receivers", 
        "senders_vs_receivers.png",
        [
            "OBSERVATION:",
            "We separated 'Broadcasters' (Senders) from 'Authorities' (Receivers).",
            "",
            "Top Senders:",
            "- High Out-Degree",
            "- Likely automated systems or announcements.",
            "",
            "Top Receivers:",
            "- High In-Degree",
            "- Likely Managers or Helpdesk."
        ]
    )

    # ==========================================
    # SLIDE 4: PageRank
    # ==========================================
    add_image_slide(
        "Influence Analysis (PageRank)",
        "pagerank_scores.png",
        [
            "WHAT IS PAGERANK?",
            "It measures influence, not just volume. A node is important if other important nodes email it.",
            "",
            "RESULTS:",
            "- The top influencers are often different from the top spammers.",
            "- These nodes represent the 'Social Capital' of the network."
        ]
    )

    # ==========================================
    # SLIDE 5: Communities
    # ==========================================
    add_image_slide(
        "Community Detection",
        "communities.png",
        [
            "WHAT IS A COMMUNITY?",
            "A group of nodes densely connected internally but sparsely connected to the outside.",
            "",
            "CONTEXT:",
            "These likely represent departments, project teams, or research labs working together."
        ]
    )

    # ==========================================
    # SLIDE 6: Resilience (Text Only)
    # ==========================================
    resilience_points = [
        "Hypothesis: Does the network break if the Top Sender (Hub) is removed?",
        " ",
        "Experiment:",
        "1. Measured size of the Core (Largest Strongly Connected Component).",
        "2. Removed Node #0 (Top Sender).",
        "3. Recalculated Core size.",
        " ",
        "Result:",
        "The Core size remained stable.",
        " ",
        "Conclusion:",
        "The network is ROBUST. Communication relies on multiple pathways, not just one central figure."
    ]
    add_bullet_slide("Network Resilience Test", resilience_points)

    # ==========================================
    # SLIDE 7: Summary
    # ==========================================
    summary_points = [
        "The network is 'Scale-Free': A few hubs dominate, but the majority have few connections.",
        "Clear Hierarchy: Distinct separation between information sources (Senders) and sinks (Receivers).",
        "Robustness: The institution's communication flow survives the loss of key players.",
        "Communities: Distinct subgroups exist, reflecting the departmental structure of the organization."
    ]
    add_bullet_slide("Conclusion", summary_points)

    # Save the file
    filename = "EU_Email_Network_Analysis.pptx"
    prs.save(filename)
    print(f"Presentation saved successfully as '{filename}'")
    print("You can open this file directly in LibreOffice Impress.")

if __name__ == "__main__":
    create_presentation()